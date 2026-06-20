from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set widescreen 16:9 (roughly 13.33 in x 7.5 in)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colors
    DARK_GREEN = RGBColor(0, 77, 64)   # #004D40
    ORANGE = RGBColor(255, 140, 0)     # #FF8C00
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(240, 240, 240)

    def add_flat_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_to_slide(slide, title_text, color=DARK_GREEN):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = color

    # Slide 1: Capa
    slide_layout = prs.slide_layouts[6] # Empty layout
    slide = prs.slides.add_slide(slide_layout)
    add_flat_background(slide, DARK_GREEN)

    # Background shape for style
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(4.5), Inches(2), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Prestação de Contas e Gestão Operacional"
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = WHITE
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Lagoa Jardins"
    p2.font.size = Pt(32)
    p2.font.color.rgb = ORANGE
    
    # Presenter
    pres_box = slide.shapes.add_textbox(Inches(4.5), Inches(5.5), Inches(8), Inches(1))
    tf_p = pres_box.text_frame
    p3 = tf_p.paragraphs[0]
    p3.text = "Jonatan Almeida - Gerente Geral"
    p3.font.size = Pt(24)
    p3.font.color.rgb = WHITE

    # Slide 2: Indicadores (KPIs)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_flat_background(slide, WHITE)
    add_title_to_slide(slide, "Nossos Indicadores (KPIs) e Resultados")

    # Metrics layout (3 columns)
    metrics = [
        ("GRI", "Reputação Online", "95.0 (+8.01)", "Março"),
        ("CSAT", "Satisfação Unidade", "4.04 (+0.13)", "Março"),
        ("NPS", "Lealdade", "51.85 (+17.07)", "Março"),
    ]
    
    for i, (m_title, desc, val, month) in enumerate(metrics):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i*4.2), Inches(1.5), Inches(3.8), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_GREEN
        box.line.color.rgb = ORANGE
        
        # Text in box
        tf = slide.shapes.add_textbox(Inches(0.6 + i*4.2), Inches(1.6), Inches(3.6), Inches(2.3)).text_frame
        p = tf.paragraphs[0]
        p.text = m_title
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = WHITE
        
        p_val = tf.add_paragraph()
        p_val.text = val
        p_val.font.bold = True
        p_val.font.size = Pt(36)
        p_val.font.color.rgb = ORANGE
        
    # April Partial
    apr_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(1.5))
    tf = apr_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Abril (Parcial): Evolução Contínua"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GREEN
    
    p_sub = tf.add_paragraph()
    p_sub.text = "GRI: 96.2  |  CSAT: 4.15  |  NPS: 46.15"
    p_sub.font.size = Pt(32)
    p_sub.font.color.rgb = ORANGE

    # Slide 3: Gestão Financeira
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_flat_background(slide, WHITE)
    add_title_to_slide(slide, "Gestão Financeira (Taxa Condominial)")
    
    # Table like layout
    items = [
        ("Casas 2 Quartos", "R$ 88,95", "R$ 84,50 (até dia 5)"),
        ("Casas 3 Quartos", "R$ 107,35", "R$ 101,98 (até dia 5)")
    ]
    
    for i, (house, full, discounted) in enumerate(items):
        row_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5 + i*1.2), Inches(12), Inches(1))
        row_box.fill.solid()
        row_box.fill.fore_color.rgb = GRAY
        row_box.line.color.rgb = DARK_GREEN
        
        tf = slide.shapes.add_textbox(Inches(0.7), Inches(1.6 + i*1.2), Inches(11), Inches(0.8)).text_frame
        p = tf.paragraphs[0]
        p.text = f"{house}: {full} "
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GREEN
        
        run = p.add_run()
        run.text = f" » Com 5% desconto: {discounted}"
        run.font.bold = True
        run.font.color.rgb = ORANGE

    # Clube de Vantagens
    club_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(12), Inches(2))
    club_box.fill.solid()
    club_box.fill.fore_color.rgb = DARK_GREEN
    
    tf = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(11), Inches(1.5)).text_frame
    p = tf.paragraphs[0]
    p.text = "Clube de Vantagens"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = ORANGE
    
    p2 = tf.add_paragraph()
    p2.text = "• Desconto de 25% na pensão antecipada"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    
    p3 = tf.add_paragraph()
    p3.text = "• 10% de desconto na estação Yellot Mob"
    p3.font.size = Pt(22)
    p3.font.color.rgb = WHITE

    # Helper for Maintenance Slides
    def add_maintenance_slide(prs, title, text_items, placeholder_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_flat_background(slide, WHITE)
        add_title_to_slide(slide, title)
        
        # Image placeholder
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = GRAY
        placeholder.line.color.rgb = DARK_GREEN
        placeholder.line.dash_style = 2 # Dashed
        
        tf_ph = placeholder.text_frame
        p_ph = tf_ph.paragraphs[0]
        p_ph.text = f"[ESPAÇO PARA IMAGEM]\n{placeholder_text}"
        p_ph.alignment = PP_ALIGN.CENTER
        p_ph.font.size = Pt(18)
        p_ph.font.color.rgb = DARK_GREEN
        
        # Text box
        text_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(5))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for item in text_items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_GREEN
            p.space_after = Pt(12)

    # Slide 4: Gestão Hídrica
    add_maintenance_slide(prs, "Ações de Manutenção: Gestão Hídrica", 
                          ["Intervenção rápida na bomba do poço (queima resolvida).", 
                           "Impacto: Redução drástica no custo de água.", 
                           "Segurança: Laudo mensal de potabilidade garantido."], 
                          "Manutenção da bomba ou poço artesiano")

    # Slide 5: Conforto e Privacidade
    add_maintenance_slide(prs, "Ações de Manutenção: Conforto e Privacidade", 
                          ["Reforma de 13 sofás (estrutura e estofado).", 
                           "Troca de insulfilm danificados (térmico/privacidade).", 
                           "Manutenção de janelas, portas, portais e alizares."], 
                          "Sofás reformados ou insulfilm")

    # Slide 6: Prevenção e Melhorias
    add_maintenance_slide(prs, "Ações de Manutenção: Prevenção", 
                          ["Dedetização preventiva em 130 casas (período de chuvas).", 
                           "Reforma completa da laje do Restaurante.", 
                           "Impermeabilização e pintura interna (Custo ZERO ao condomínio)."], 
                          "Dedetização ou laje do restaurante")

    # Slide 7: Processos Inteligentes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_flat_background(slide, WHITE)
    add_title_to_slide(slide, "Processos Inteligentes e Agilidade")
    
    # Process Boxes
    proc_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    proc_box.fill.solid()
    proc_box.fill.fore_color.rgb = GRAY
    proc_box.line.color.rgb = ORANGE
    
    tf = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(11), Inches(4.5)).text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Check-in e Voucher"
    p1.font.bold = True
    p1.font.size = Pt(24)
    p1.font.color.rgb = DARK_GREEN
    
    p1s = tf.add_paragraph()
    p1s.text = "Marque com antecedência. Voucher obrigatório no check-in."
    p1s.font.size = Pt(18)
    
    p2 = tf.add_paragraph()
    p2.text = "FNRH Online (Lei 15.004/2025)"
    p2.font.bold = True
    p2.font.size = Pt(24)
    p2.font.color.rgb = ORANGE
    p2.space_before = Pt(18)
    
    p2s = tf.add_paragraph()
    p2s.text = "Link enviado 5 dias antes. Extingue a burocracia na recepção.\nLogin: Data Check-in + (Localizador/Sobrenome OU CPF)."
    p2s.font.size = Pt(18)

    # Slide 8: Segurança Jurídica (ECA)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_flat_background(slide, RGBColor(183, 28, 28)) # Warning Red background for legal 
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⚠ Segurança Jurídica: Regras para Menores (ECA)"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
    tf = text_box.text_frame
    
    def add_warning(tf, title, detail):
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        ps = tf.add_paragraph()
        ps.text = detail
        ps.font.size = Pt(18)
        ps.font.color.rgb = WHITE
        ps.space_after = Pt(10)
        
    add_warning(tf, "Com os Pais:", "Indispensável RG ou Certidão + Documento dos pais.")
    add_warning(tf, "Com Terceiros:", "Autorização por escrito com FIRMA RECONHECIDA.")
    add_warning(tf, "Desacompanhado:", "Autorização judicial ou dos pais (firma reconhecida).")

    # Slide 9: Recreação
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_flat_background(slide, ORANGE)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ritmo, Cor e Diversão: Feriadão na Lagoa"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = DARK_GREEN
    
    acts_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf = acts_box.text_frame
    items = ["Alongamento aquático", "Rádio Jardins", "Desafio do Balão", "Quiz Maluco", "Oficina de Slime"]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"★ {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        
    # Easter Special
    easter_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5), Inches(4))
    easter_box.fill.solid()
    easter_box.fill.fore_color.rgb = WHITE
    
    tf_e = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(4.6), Inches(3.6)).text_frame
    p = tf_e.paragraphs[0]
    p.text = "Especial de Páscoa"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = ORANGE
    
    for act in ["Coelho Colorido (04/04)", "Master Chef (04/04)", "Entrega de Ovos (05/04)"]:
        pa = tf_e.add_paragraph()
        pa.text = f"• {act}"
        pa.font.size = Pt(20)
        pa.font.color.rgb = DARK_GREEN

    # Slide 10: Agenda e Contatos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_flat_background(slide, DARK_GREEN)
    
    add_title_to_slide(slide, "Agenda Musical e Contatos", color=WHITE)
    
    # Music
    music_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(4))
    tf_m = music_box.text_frame
    p_m = tf_m.paragraphs[0]
    p_m.text = "Shows Confirmados (18h):"
    p_m.font.bold = True
    p_m.font.size = Pt(24)
    p_m.font.color.rgb = ORANGE
    
    shows = ["04/04: D'Borges", "11/04: Talles Guimarães", "18/04: Neto Alves", "25/04: Marcel Araújo"]
    for show in shows:
        ps = tf_m.add_paragraph()
        ps.text = show
        ps.font.size = Pt(20)
        ps.font.color.rgb = WHITE
        
    # Contacts
    cont_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(6), Inches(4))
    tf_c = cont_box.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.text = "Canais de Atendimento:"
    p_c.font.bold = True
    p_c.font.size = Pt(24)
    p_c.font.color.rgb = ORANGE
    
    conts = ["Central: 0800 960 5040", "Recepção: (64) 3513-1100"]
    for cont in conts:
        pc = tf_c.add_paragraph()
        pc.text = cont
        pc.font.size = Pt(20)
        pc.font.color.rgb = WHITE

    # Save
    prs.save("Apresentacao_Lagoa_Jardins_Aprimorada.pptx")
    print("Presentation created successfully: Apresentacao_Lagoa_Jardins_Aprimorada.pptx")

if __name__ == "__main__":
    create_presentation()
