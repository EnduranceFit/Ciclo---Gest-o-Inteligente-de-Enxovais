import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Image Paths from current brain directory
BRAIN_DIR = r"C:\Users\Jonat\.gemini\antigravity\brain\11762cc6-0706-4b33-a5cc-0bc6b1ed1f18"
COVER_IMG = os.path.join(BRAIN_DIR, "cover_background_jardins_1776880975375.png")
RECREATION_IMG = os.path.join(BRAIN_DIR, "recreation_background_1776881057410.png")

def create_premium_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Professional Palette
    DARK_GREEN = RGBColor(0, 51, 42)    # #00332A - Master deep green
    ACCENT_ORANGE = RGBColor(255, 102, 0) # #FF6600 - Sophisticated Orange
    WHITE = RGBColor(255, 255, 255)
    OFF_WHITE = RGBColor(248, 249, 250)
    SOFT_GRAY = RGBColor(224, 224, 224)

    def set_bg_color(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_sidebar(slide):
        # Elegant sidebar on the left
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_GREEN
        shape.line.fill.background()
        
        # Accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0), Inches(0.05), Inches(7.5))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_ORANGE
        line.line.fill.background()

    def add_title(slide, text, subtitle=None):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = text.upper()
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = DARK_GREEN
        p.font.name = "Calibri Light"
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(14)
            p2.font.color.rgb = ACCENT_ORANGE
            p2.font.name = "Calibri"

    # Slide 1: Cover (Premium Cinematográfico)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if os.path.exists(COVER_IMG):
        slide.shapes.add_picture(COVER_IMG, Inches(0), Inches(0), height=Inches(7.5))
    else:
        set_bg_color(slide, DARK_GREEN)

    # Dark overlay for better text readability
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.fill.transparency = 0.4
    overlay.line.fill.background()

    # Title box (Glassmorphism look simulation)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_GREEN
    box.fill.transparency = 0.2
    box.line.color.rgb = ACCENT_ORANGE
    box.line.width = Pt(2)

    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = "PRESTAÇÃO DE CONTAS"
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "LAGOA JARDINS"
    p2.font.bold = True
    p2.font.size = Pt(32)
    p2.font.color.rgb = ACCENT_ORANGE
    
    p3 = tf.add_paragraph()
    p3.text = "\nGestão Operacional | Jonatan Almeida"
    p3.font.size = Pt(20)
    p3.font.color.rgb = WHITE

    # Slide 2: KPIs (Dashboard Style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, OFF_WHITE)
    add_sidebar(slide)
    add_title(slide, "Nossos Indicadores e Resultados", "Acompanhamento Mensal (Março/Abril)")

    metrics = [
        ("GRI", "Reputação Online", "95.0", "+8.01", "📈"),
        ("CSAT", "Satisfação Unidade", "4.04", "+0.13", "👥"),
        ("NPS", "Lealdade", "51.85", "+17.07", "🏆")
    ]

    for i, (name, label, val, trend, icon) in enumerate(metrics):
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1 + i*4), Inches(2), Inches(3.5), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = SOFT_GRAY
        
        # Content
        tf = slide.shapes.add_textbox(Inches(1.2 + i*4), Inches(2.2), Inches(3.1), Inches(3.1)).text_frame
        p_icon = tf.paragraphs[0]
        p_icon.text = icon
        p_icon.font.size = Pt(24)
        
        p_name = tf.add_paragraph()
        p_name.text = name
        p_name.font.bold = True
        p_name.font.size = Pt(28)
        p_name.font.color.rgb = DARK_GREEN
        
        p_val = tf.add_paragraph()
        p_val.text = val
        p_val.font.bold = True
        p_val.font.size = Pt(48)
        p_val.font.color.rgb = DARK_GREEN
        
        p_trend = tf.add_paragraph()
        p_trend.text = f"Tendência: {trend}"
        p_trend.font.size = Pt(16)
        p_trend.font.color.rgb = RGBColor(0, 150, 0) # Green for positive

    # Slide 3: Financeiro
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, OFF_WHITE)
    add_sidebar(slide)
    add_title(slide, "Gestão Financeira", "Taxa Condominial Semanal")

    pricing = [
        ("Casas 2 Quartos", "R$ 88,95", "R$ 84,50 (5% Desc.)"),
        ("Casas 3 Quartos", "R$ 107,35", "R$ 101,98 (5% Desc.)")
    ]

    for i, (type, full, disc) in enumerate(pricing):
        # Horizontal Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2 + i*1.5), Inches(11.5), Inches(1.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = WHITE
        bar.line.color.rgb = DARK_GREEN
        
        tf = slide.shapes.add_textbox(Inches(1.5), Inches(2.1 + i*1.5), Inches(10.5), Inches(1)).text_frame
        p = tf.paragraphs[0]
        p.text = f"{type}:   "
        p.font.size = Pt(24)
        run = p.add_run()
        run.text = full
        run.font.bold = True
        
        p2 = tf.add_paragraph()
        p2.text = f"Valor com desconto até dia 05: {disc}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = ACCENT_ORANGE

    # Clube Vantagens banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.5), Inches(11.5), Inches(1.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DARK_GREEN
    
    tf = banner.text_frame
    p = tf.add_paragraph()
    p.text = "👑 CLUBE DE VANTAGENS"
    p.font.bold = True
    p.font.size = Pt(26)
    p.font.color.rgb = ACCENT_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "25% Off Pensão Antecipada | 10% Off Estação Yellot Mob"
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # Helper for Professional Content Slide
    def add_content_slide(prs, title, subtitle, bullets, img_placeholder_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg_color(slide, OFF_WHITE)
        add_sidebar(slide)
        add_title(slide, title, subtitle)
        
        # Left Content Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(6), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = SOFT_GRAY
        
        tf = slide.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(5.4), Inches(4)).text_frame
        tf.word_wrap = True
        for b in bullets:
            p = tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_GREEN
            p.space_after = Pt(14)
            
        # Right Image Frame
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(2), Inches(5.3), Inches(4.5))
        frame.fill.solid()
        frame.fill.fore_color.rgb = SOFT_GRAY
        frame.line.color.rgb = ACCENT_ORANGE
        frame.line.dash_style = 2
        
        tf_f = frame.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"FOTO: {img_placeholder_text}"
        p_f.alignment = PP_ALIGN.CENTER
        p_f.font.size = Pt(18)
        p_f.font.color.rgb = DARK_GREEN

    # Slide 4: Hidrica
    add_content_slide(prs, "Ações de Manutenção", "Gestão Hídrica Inteligente", 
                      ["Intervenção imediata na bomba do poço (queima resolvida).", 
                       "Poço artesiano como fonte principal = ECONOMIA REAL.", 
                       "Potabilidade Garantida: Laudo mensal rigoroso."], "Bomba / Poço")

    # Slide 5: Conforto
    add_content_slide(prs, "Ações de Manutenção", "Conforto e Privacidade", 
                      ["Reforma estrutural e técnica de 13 sofás.", 
                       "Troca de insulfilm danificado nas janelas.", 
                       "Revisão de aberturas: Janelas, Portas e Alizares."], "Sofás / Insulfilm")

    # Slide 6: Prevenção
    add_content_slide(prs, "Ações de Manutenção", "Prevenção e Melhorias", 
                      ["Dedetização preventiva em 100% das 130 casas.", 
                       "Reforma da laje e impermeabilização do Restaurante.", 
                       "IMPORTANTE: Obra realizada SEM CUSTO ao condomínio."], "Restaurant / Dedetização")

    # Slide 7: Processos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, OFF_WHITE)
    add_sidebar(slide)
    add_title(slide, "Processos Inteligentes", "Check-in e FNRH Digital")
    
    # Process Row
    for i, (title, steps) in enumerate([
        ("Check-in & Voucher", "Voucher enviado na reserva. Apresentação obrigatória para segurança."),
        ("FNRH Online (Lei 15.004/2025)", "Link enviado 5 dias antes. Preenchimento não editável.\nAgilidade total na chegada.")
    ]):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1 + i*6), Inches(2), Inches(5.5), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = DARK_GREEN
        
        tf = slide.shapes.add_textbox(Inches(1.2 + i*6), Inches(2.2), Inches(5.1), Inches(4)).text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_ORANGE
        
        pd = tf.add_paragraph()
        pd.text = f"\n{steps}"
        pd.font.size = Pt(20)
        pd.font.color.rgb = DARK_GREEN

    # Slide 8: ECA (Redesign Professional Alerta)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, DARK_GREEN)
    
    # Alerta icon shape
    shape = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_PROCESS, Inches(4), Inches(1), Inches(5.33), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_ORANGE
    
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = "SIRIUS: SEGURANÇA JURÍDICA (ECA)"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Legal points
    points = [
        ("Pais", "RG/Certidão + Doc Original dos Pais"),
        ("Terceiros", "Autorização reconhecida em CARTÓRIO"),
        ("Sozinhos", "Autorização Judicial ou Parental (Firma reconhecida)")
    ]
    
    for i, (who, rule) in enumerate(points):
        row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3 + i*1.2), Inches(11.33), Inches(1))
        row.fill.solid()
        row.fill.fore_color.rgb = RGBColor(0, 64, 53)
        row.line.color.rgb = ACCENT_ORANGE
        
        tf = slide.shapes.add_textbox(Inches(1.5), Inches(3.1 + i*1.2), Inches(10), Inches(0.8)).text_frame
        p = tf.paragraphs[0]
        p.text = f"{who}:  "
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_ORANGE
        run = p.add_run()
        run.text = rule
        run.font.bold = False
        run.font.color.rgb = WHITE

    # Slide 9: Recreação (Vibrant IA Background)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if os.path.exists(RECREATION_IMG):
        slide.shapes.add_picture(RECREATION_IMG, Inches(0), Inches(0), height=Inches(7.5))
    
    add_title(slide, "Ritmo, Cor e Diversão", "Programação Especial de Páscoa")
    
    # Overly to pop text
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2), Inches(6), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_GREEN
    box.fill.transparency = 0.1
    box.line.color.rgb = ACCENT_ORANGE
    
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = "Destaques Feriadão:"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_ORANGE
    
    items = ["Coelho Colorido (04/04)", "Master Chef Junior (04/04)", "Entrega de Ovos (05/04)", "Rádio Jardins & Gincanas"]
    for item in items:
        pi = tf.add_paragraph()
        pi.text = f"★ {item}"
        pi.font.size = Pt(22)
        pi.font.color.rgb = WHITE
        pi.space_before = Pt(10)

    # Slide 10: Encerramento
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sidebar(slide)
    set_bg_color(slide, DARK_GREEN)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AGENDA MUSICAL & CONTATOS"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = ACCENT_ORANGE
    
    # Music list
    music = ["04/04 • D'Borges", "11/04 • Talles Guimarães", "18/04 • Neto Alves", "25/04 • Marcel Araújo"]
    for i, m in enumerate(music):
        p = tf.add_paragraph()
        p.text = m
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        
    # Bottom Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(6), Inches(12.93), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 31, 26)
    
    tf_b = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(1)).text_frame
    p = tf_b.paragraphs[0]
    p.text = "Central: 0800 960 5040  |  Recepção: (64) 3513-1100"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_ORANGE

    # Save
    name = "Apresentacao_Lagoa_Jardins_PREMIUM.pptx"
    prs.save(name)
    print(f"Professional presentation saved: {name}")

if __name__ == "__main__":
    create_premium_presentation()
